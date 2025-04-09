import streamlit as st
import google.generativeai as genai
import os
import tempfile
from bs4 import BeautifulSoup
import requests
import git
import shutil

# 🔐 Gemini API Setup
api_key = st.secrets["GEMINI_API_KEY"]
genai.configure(api_key=api_key)
model = genai.GenerativeModel("gemini-2.0-flash")

st.set_page_config(page_title="AI Search + QA Agents", layout="wide",page_icon="🔎")
st.title("🔎 Search + QA Agents")
st.markdown("Effortlessly query documents, websites, or GitHub repos with Gemini + RAG ")

option = st.sidebar.selectbox("Choose Agent", [
    "Multi-Doc QA Bot with RAG",
    "Website Chat Agent (RAG from URL)",
    "GitHub Repo Assistant"
])

import os
import tempfile
import pandas as pd
import fitz  # PyMuPDF
import docx2txt
import json
import pptx
import openpyxl
from bs4 import BeautifulSoup
import ebooklib
from ebooklib import epub

# --- Multi-Doc QA Agent ---
if option == "Multi-Doc QA Bot with RAG":
    st.subheader("📄 Multi-Document QA with RAG")
    uploaded_files = st.file_uploader(
        "Upload multiple documents",
        accept_multiple_files=True,
        type=["pdf", "txt", "docx", "csv", "json", "md", "pptx", "xlsx", "html", "epub"]
    )
    question = st.text_input("Ask a question about the documents")

    if st.button("Get Answer") and uploaded_files and question:
        combined_text = ""
        for file in uploaded_files:
            filename = file.name.lower()
            if filename.endswith(".txt") or filename.endswith(".md"):
                combined_text += file.read().decode("utf-8") + "\n"

            elif filename.endswith(".pdf"):
                with fitz.open(stream=file.read(), filetype="pdf") as doc:
                    for page in doc:
                        combined_text += page.get_text()

            elif filename.endswith(".docx"):
                temp_path = os.path.join(tempfile.gettempdir(), file.name)
                with open(temp_path, "wb") as f:
                    f.write(file.read())
                combined_text += docx2txt.process(temp_path)

            elif filename.endswith(".csv"):
                df = pd.read_csv(file)
                combined_text += df.to_string(index=False)

            elif filename.endswith(".json"):
                data = json.load(file)
                combined_text += json.dumps(data, indent=2)

            elif filename.endswith(".pptx"):
                prs = pptx.Presentation(file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            combined_text += shape.text + "\n"

            elif filename.endswith(".xlsx"):
                df = pd.read_excel(file, sheet_name=None)
                for name, sheet in df.items():
                    combined_text += f"\nSheet: {name}\n" + sheet.to_string(index=False)

            elif filename.endswith(".html"):
                soup = BeautifulSoup(file.read(), "html.parser")
                combined_text += soup.get_text()

            elif filename.endswith(".epub"):
                book = epub.read_epub(file)
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        soup = BeautifulSoup(item.get_content(), "html.parser")
                        combined_text += soup.get_text()

        prompt = f"""You are an AI assistant. Use the context below to answer the question:

Context:
{combined_text}

Question: {question}
"""
        response = model.generate_content(prompt)
        st.write(response.text)


# --- Website Chat Agent ---
elif option == "Website Chat Agent (RAG from URL)":
    st.subheader("🌐 Website QA Agent")
    url = st.text_input("Enter website URL")
    question = st.text_input("What do you want to know?")

    if st.button("Ask Website") and url and question:
        try:
            r = requests.get(url)
            soup = BeautifulSoup(r.text, "html.parser")
            content = soup.get_text()
            prompt = f"""Use the following website content to answer:

Website Content:
{content}

Question: {question}
"""
            response = model.generate_content(prompt)
            st.write(response.text)
        except Exception as e:
            st.error(f"Failed to load content: {e}")

# --- GitHub Repo Assistant ---
elif option == "GitHub Repo Assistant":
    st.subheader("🐙 GitHub Repo QA")
    repo_url = st.text_input("Enter GitHub Repository URL")
    question = st.text_input("Ask something about the repo")

    if st.button("Ask Repo") and repo_url and question:
        try:
            temp_dir = tempfile.mkdtemp()
            git.Repo.clone_from(repo_url, temp_dir)
            combined_code = ""
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    if file.endswith((".py", ".md", ".txt", ".js", ".java")):
                        with open(os.path.join(root, file), "r", errors='ignore') as f:
                            combined_code += f.read() + "\n"

            prompt = f"""You are an AI software assistant. Use the following codebase and README to answer:

Code + Docs:
{combined_code[:20000]}

Question: {question}
"""
            response = model.generate_content(prompt)
            st.write(response.text)
            shutil.rmtree(temp_dir)
        except Exception as e:
            st.error(f"Error cloning or reading the repo: {e}")
